#!/usr/bin/env python3
# notes_extractor.py
# Extract speaker notes from a .pptx into a plain .txt file (with slide/page markers).

import argparse
from pathlib import Path
from pptx import Presentation


def extract_notes(pptx_path: Path) -> list[tuple[int, str]]:
    prs = Presentation(str(pptx_path))
    out: list[tuple[int, str]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes = ""

        # Normalize line endings a bit
        notes = notes.replace("\r\n", "\n").replace("\r", "\n").strip()

        out.append((idx, notes))
    return out


def default_output_path(input_path: Path) -> Path:
    # input.pptx -> input_notes.txt
    return input_path.with_name(f"{input_path.stem}_notes.txt")


def main():
    parser = argparse.ArgumentParser(
        description="Extract PowerPoint speaker notes into a plain text file."
    )
    parser.add_argument("-i", "--input", required=True, help="Input .pptx file path")
    parser.add_argument(
        "-o",
        "--output",
        default=None,
        help="Output .txt file path (default: <input_stem>_notes.txt next to input)",
    )
    args = parser.parse_args()

    in_path = Path(args.input).expanduser().resolve()
    if not in_path.exists():
        raise SystemExit(f"[ERROR] Input file not found: {in_path}")
    if in_path.suffix.lower() != ".pptx":
        raise SystemExit("[ERROR] Input must be a .pptx file (not .ppt, .pptm).")

    out_path = (
        Path(args.output).expanduser().resolve()
        if args.output
        else default_output_path(in_path)
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)

    notes_by_slide = extract_notes(in_path)

    with out_path.open("w", encoding="utf-8") as f:
        f.write(f"# Speaker Notes Extract\n")
        f.write(f"# Input: {in_path.name}\n\n")

        for slide_no, notes in notes_by_slide:
            f.write(f"===== Page {slide_no} =====\n")
            if notes:
                f.write(notes + "\n")
            else:
                f.write("(no notes)\n")
            f.write("\n")

    print(f"[OK] Saved: {out_path}")


if __name__ == "__main__":
    main()



